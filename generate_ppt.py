#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
日本跨境电商 ERP 系统 - 合作伙伴版演示文档（优化版）
生成日期：2026-03-07
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_ppt():
    prs = Presentation()
    
    # 设置宽屏 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 配色方案
    BLUE = RGBColor(41, 98, 255)      # 科技蓝
    ORANGE = RGBColor(255, 127, 14)   # 活力橙
    DARK = RGBColor(44, 62, 80)       # 深灰
    LIGHT = RGBColor(149, 165, 166)   # 浅灰
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = LIGHT
        return slide
    
    def add_content_slide(title, content_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
        
        body = slide.placeholders[1].text_frame
        body.clear()
        
        for i, line in enumerate(content_lines):
            if i == 0:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.space_after = Pt(12)
            if line.startswith("•") or line.startswith("→"):
                p.level = 0
            elif line.startswith("  -"):
                p.level = 1
        return slide
    
    # ========== 第 1 页：封面 ==========
    slide = add_title_slide(
        "日本跨境电商 ERP 系统",
        "一站式自动化解决方案 · 助力卖家效率提升 10 倍"
    )
    
    # ========== 第 2 页：演示信息 ==========
    slide = add_title_slide(
        "合作伙伴版演示文档",
        "2026 年 3 月 · 二狗子 🐕"
    )
    
    # ========== 第 3 页：目录 ==========
    add_content_slide("📋 目录", [
        "一、项目背景与机会",
        "二、产品定位",
        "三、核心功能",
        "四、商业价值",
        "五、投入与回报",
        "六、合作模式",
        "七、下一步行动"
    ])
    
    # ========== 第 4 页：项目背景 ==========
    add_content_slide("一、项目背景与机会", [
        "🎯 市场机会",
        "• 日本电商市场：¥20 万亿+/年",
        "• 跨境卖家：10 万+",
        "• 核心痛点：多平台管理复杂、人工效率低、合规风险高",
        "",
        "❌ 现状挑战",
        "• 卖家平均运营 3-5 个店铺（亚马逊、乐天、雅虎）",
        "• 订单/库存/物流/财务全靠人工",
        "• 错误率高、无法规模化",
        "",
        "✅ 解决方案",
        "• 一个系统管理所有店铺",
        "• 自动化处理 90% 日常操作",
        "• 日本本地化合规支持"
    ])
    
    # ========== 第 5 页：产品定位 ==========
    add_content_slide("二、产品定位", [
        "🏷️ 一句话说明",
        "→ 面向日本市场的跨境电商管理系统",
        "→ 自动处理订单、库存、物流、财务全流程",
        "",
        "👥 目标客户",
        "• 中小卖家（1-10 人）：人手不足，效率低",
        "• 中型卖家（10-50 人）：多店管理混乱，数据不统一",
        "• 大型卖家（50 人+）：需要系统化，合规要求高",
        "",
        "🎖️ 核心优势",
        "• 日本本地化：税务、物流、语言全支持",
        "• 多平台对接：亚马逊、乐天、雅虎一键接入",
        "• AI 智能助手：自动文案、智能客服、经营分析",
        "• SaaS 模式：按需订阅，快速上线"
    ])
    
    # ========== 第 6 页：核心功能（上）- 业务模块 ==========
    add_content_slide("三、核心功能 - 业务模块", [
        "📦 8 大核心业务模块",
        "",
        "店铺管理 → 一个后台管所有店",
        "订单中心 → 30 分钟→3 分钟（10 倍提升）",
        "库存管理 → 防止超卖，准确率 99.9%",
        "物流管理 → 对接 Yamato、Sagawa，运费 -15%",
        "财务管理 → 3 天→3 小时（24 倍提升）",
        "客服管理 → 智能回复，效率 3 倍",
        "BI 分析 → 数据驱动决策",
        "AI 助手 → 文案/简报/问答，效率 +50%"
    ])
    
    # ========== 第 7 页：核心功能（下）- 支撑模块 ==========
    add_content_slide("三、核心功能 - 支撑模块", [
        "🛠️ 7 大支撑模块",
        "",
        "• 商品中心：一次录入，多平台分发（效率 5 倍）",
        "• 采购管理：智能补货，断货率 -80%",
        "• 广告管理：多平台数据统一，ROI +20%",
        "• 规则引擎：自动化流程，减少 90% 重复操作",
        "• 合规管理：日本税务/法律合规检查",
        "• 权限管理：多角色权限控制",
        "• 系统管理：监控告警，稳定性 99.9%"
    ])
    
    # ========== 第 8 页：商业价值 ==========
    add_content_slide("四、商业价值", [
        "💰 客户收益（年销 5000 万日元卖家）",
        "",
        "效率提升：",
        "• 订单处理：10 倍 | 商品上新：6 倍",
        "• 财务结算：24 倍 | 客服响应：12 倍",
        "",
        "成本降低：",
        "• 人工成本：-50% | 物流成本：-15%",
        "• 库存成本：-30% | 错误成本：-90%",
        "",
        "收入增长：",
        "• 店铺数量：3 个→10 个",
        "• 转化率：+15% | 广告 ROI：+20%",
        "",
        "📊 投资回报率：44 倍！",
        "投入¥10 万/年 → 净收益¥440 万/年"
    ])
    
    # ========== 第 9 页：投入与回报 ==========
    add_content_slide("五、投入与回报", [
        "🛠️ 开发投入：¥165 万",
        "• 人力成本：¥150 万（3 后端 +2 前端 +1 测试 +1 产品 +1UI）",
        "• 服务器/工具：¥10 万/年",
        "• 第三方 API：¥5 万/年",
        "",
        "📈 收益预测（SaaS 订阅）",
        "• 基础版：100 客户 × ¥5 万 = ¥500 万",
        "• 专业版：50 客户 × ¥15 万 = ¥750 万",
        "• 企业版：10 客户 × ¥50 万 = ¥500 万",
        "",
        "💵 合计：160 客户，¥1750 万/年",
        "毛利率：70%（SaaS 模式）"
    ])
    
    # ========== 第 10 页：盈利时间线 ==========
    add_content_slide("五、投入与回报 - 盈利时间线", [
        "💵 投资回收期",
        "",
        "第 1 年 → 收回成本 + 盈利",
        "第 2 年 → 净利润 ¥1500 万+",
        "第 3 年 → 累计净利润 ¥3000 万+",
        "",
        "🎯 关键假设",
        "• 获客成本：¥5 万/客户",
        "• 客户留存率：80%/年",
        "• 毛利率：70%"
    ])
    
    # ========== 第 11 页：合作模式 ==========
    add_content_slide("六、合作模式", [
        "🤝 三种合作方式",
        "",
        "方案 A：联合开发（推荐）",
        "• 我方：技术 + 产品 | 合作方：市场 + 客户",
        "• 收益分成：50% : 50%",
        "",
        "方案 B：技术授权",
        "• 授权费：¥100 万一次性 + 10% 年营收分成",
        "• 合作方：自主运营 + 品牌定制",
        "",
        "方案 C：合资公司",
        "• 共同成立公司运营产品",
        "• 股权比例：协商 | 长期收益共享"
    ])
    
    # ========== 第 12 页：时间计划 ==========
    add_content_slide("六、合作模式 - 14 周上线计划", [
        "📅 开发时间线",
        "",
        "第 1 周    → 需求确认（确定功能范围）",
        "第 2-3 周  → 系统设计（完成架构设计）",
        "第 4-12 周 → 开发实施（完成核心功能）",
        "第 13-14 周→ 测试上线（正式发布）",
        "第 15 周起 → 市场推广（获客运营）",
        "",
        "🚀 快速启动，14 周即可上线盈利！"
    ])
    
    # ========== 第 13 页：下一步行动 ==========
    add_content_slide("七、下一步行动", [
        "🎯 立即行动",
        "",
        "本周   → 确认合作意向和模式",
        "下周   → 签订合作协议",
        "第 3 周 → 启动项目开发",
        "第 14 周→ 产品上线 🚀",
        "第 15 周→ 开始获客 💰",
        "",
        "👉 本周是关键！尽早确认，尽早收益"
    ])
    
    # ========== 第 14 页：联系我们 ==========
    slide = add_title_slide(
        "📞 联系我们",
        "随时沟通，共创双赢！\n\n制作人：二狗子 🐕\n版本：1.0 | 2026 年 3 月 7 日"
    )
    
    # 保存文件
    output_path = "/home/admin/.openclaw/workspace/日本跨境电商 ERP 系统_合作伙伴版_v2.pptx"
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
